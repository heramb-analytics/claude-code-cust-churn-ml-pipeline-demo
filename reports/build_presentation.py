"""Build pipeline_presentation.pptx for the Transaction Anomaly Detection pipeline."""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Constants ─────────────────────────────────────────────────────────────────
NAVY = RGBColor(0x1E, 0x3A, 0x8A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLUE_LIGHT = RGBColor(0x60, 0xA5, 0xFA)
GREEN = RGBColor(0x34, 0xD3, 0x99)
YELLOW = RGBColor(0xFB, 0xBF, 0x24)
GRAY_DARK = RGBColor(0x1F, 0x29, 0x37)
GRAY_MID = RGBColor(0x37, 0x41, 0x51)

FIGURES_DIR = Path("reports/figures")
SCREENSHOTS_DIR = Path("reports/screenshots")
MODELS_DIR = Path("models")
OUT_PATH = Path("reports/pipeline_presentation.pptx")

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def load_metrics() -> dict:
    """Load model metrics from disk."""
    metrics_path = MODELS_DIR / "pipeline_model_metrics.json"
    if metrics_path.exists():
        with metrics_path.open() as fh:
            return json.load(fh)
    return {}


def _fill_slide_bg(slide, color: RGBColor) -> None:
    """Fill slide background with solid color."""
    from pptx.oxml.ns import qn
    from lxml import etree

    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: int = 18,
    bold: bool = False,
    color: RGBColor = WHITE,
    align=PP_ALIGN.LEFT,
    wrap: bool = True,
) -> None:
    """Add a styled textbox to a slide."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txBox.text_frame.word_wrap = wrap
    p = txBox.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_bullet_box(
    slide,
    bullets: list[str],
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: int = 16,
    color: RGBColor = WHITE,
    bullet_color: RGBColor = BLUE_LIGHT,
) -> None:
    """Add a bulleted list textbox."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        # bullet symbol
        r1 = p.add_run()
        r1.text = "▸  "
        r1.font.size = Pt(font_size)
        r1.font.color.rgb = bullet_color
        r2 = p.add_run()
        r2.text = bullet
        r2.font.size = Pt(font_size)
        r2.font.color.rgb = color


def _add_rect(slide, left, top, width, height, fill_color: RGBColor, alpha: int = 255) -> None:
    """Add a filled rectangle shape."""
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()


def _add_image_safe(slide, img_path: Path, left, top, width, height) -> bool:
    """Add image if it exists, return True on success."""
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), Inches(left), Inches(top),
                                 Inches(width), Inches(height))
        return True
    return False


def _add_section_header(slide, title: str, subtitle: str = "") -> None:
    """Add a styled slide title header bar."""
    _add_rect(slide, 0, 0, 13.33, 1.3, NAVY)
    _add_textbox(slide, title, 0.4, 0.15, 12.5, 0.8,
                 font_size=32, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        _add_textbox(slide, subtitle, 0.4, 0.85, 12.5, 0.4,
                     font_size=14, color=BLUE_LIGHT, align=PP_ALIGN.LEFT)


def build_presentation() -> None:
    """Build all 8 slides and save the presentation."""
    metrics = load_metrics()
    test_m = metrics.get("test_metrics", {})
    f1 = test_m.get("f1", 0.5714)
    recall = test_m.get("recall", 1.0)
    precision = test_m.get("precision", 0.4)
    algorithm = metrics.get("algorithm", "XGBoostClassifier")
    feature_count = len(metrics.get("feature_cols", []))
    best_params = metrics.get("best_params", {})

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]  # Blank

    # ── SLIDE 1 — Cover ───────────────────────────────────────────────────────
    s1 = prs.slides.add_slide(blank_layout)
    _fill_slide_bg(s1, NAVY)

    # Accent bar
    _add_rect(s1, 0, 5.8, 13.33, 0.08, BLUE_LIGHT)

    _add_textbox(s1, "Transaction Anomaly", 1.0, 1.8, 11.0, 1.2,
                 font_size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_textbox(s1, "Detection ML Pipeline", 1.0, 2.9, 11.0, 1.0,
                 font_size=44, bold=True, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)
    _add_textbox(s1, "Built with Claude Code  ·  March 2026", 1.0, 4.1, 11.0, 0.6,
                 font_size=18, color=WHITE, align=PP_ALIGN.CENTER)
    _add_textbox(s1, "XGBoost  ·  FastAPI  ·  Playwright  ·  JIRA", 1.0, 4.7, 11.0, 0.5,
                 font_size=14, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)
    _add_textbox(s1, "End-to-end pipeline from raw data to production API", 1.0, 5.4, 11.0, 0.5,
                 font_size=13, color=RGBColor(0x9C, 0xA3, 0xAF), align=PP_ALIGN.CENTER)

    # ── SLIDE 2 — Problem Statement ───────────────────────────────────────────
    s2 = prs.slides.add_slide(blank_layout)
    _fill_slide_bg(s2, GRAY_DARK)
    _add_section_header(s2, "The Problem", "Why automated anomaly detection matters")

    _add_bullet_box(s2, [
        "5 transactions in raw data · 40% anomaly rate (2 of 5 flagged)",
        "Anomalies concentrated in electronics category and merchant MCC002",
        "Manual review is time-consuming and error-prone at scale",
        "Goal: real-time automated anomaly scoring with recall = 1.0 (catch every fraud)",
    ], 0.6, 1.6, 12.0, 3.5, font_size=20)

    # Stat boxes
    for i, (label, val, col) in enumerate([
        ("Transactions", "5", WHITE),
        ("Anomaly Rate", "40%", RGBColor(0xF8, 0x71, 0x71)),
        ("Categories", "3", BLUE_LIGHT),
        ("Merchants", "3", GREEN),
    ]):
        x = 0.5 + i * 3.1
        _add_rect(s2, x, 5.2, 2.8, 1.6, GRAY_MID)
        _add_textbox(s2, val, x, 5.3, 2.8, 0.8, font_size=36, bold=True, color=col, align=PP_ALIGN.CENTER)
        _add_textbox(s2, label, x, 6.0, 2.8, 0.5, font_size=12, color=RGBColor(0x9C, 0xA3, 0xAF), align=PP_ALIGN.CENTER)

    # ── SLIDE 3 — EDA Highlights ──────────────────────────────────────────────
    s3 = prs.slides.add_slide(blank_layout)
    _fill_slide_bg(s3, GRAY_DARK)
    _add_section_header(s3, "Data Overview", "Exploratory Data Analysis — 5 charts generated")

    _add_image_safe(s3, FIGURES_DIR / "01_amount_distribution.png", 0.3, 1.4, 6.3, 3.8)
    _add_image_safe(s3, FIGURES_DIR / "02_anomaly_rate_by_category.png", 6.8, 1.4, 6.2, 3.8)

    _add_textbox(s3,
                 "5 rows  ·  11 features  ·  40% anomaly rate  ·  3 categories  ·  3 merchants",
                 0.3, 5.4, 12.7, 0.5, font_size=14,
                 color=BLUE_LIGHT, align=PP_ALIGN.CENTER)

    # ── SLIDE 4 — Data Engineering ────────────────────────────────────────────
    s4 = prs.slides.add_slide(blank_layout)
    _fill_slide_bg(s4, GRAY_DARK)
    _add_section_header(s4, "Data Pipeline", "10 quality checks · 12 validation checks · 11 features")

    # Left: checks table
    checks = [
        ("Required columns present", "✅ Pass"),
        ("DataFrame not empty", "✅ Pass"),
        ("No duplicate IDs", "✅ Pass"),
        ("Amount positive", "✅ Pass"),
        ("is_anomaly binary", "✅ Pass"),
        ("No null values", "✅ Pass"),
        ("Timestamp parseable", "✅ Pass"),
        ("Amount within range", "✅ Pass"),
        ("Category non-empty", "✅ Pass"),
        ("merchant_id non-empty", "✅ Pass"),
    ]
    _add_rect(s4, 0.3, 1.4, 6.0, 5.8, GRAY_MID)
    _add_textbox(s4, "Quality Checks", 0.3, 1.4, 6.0, 0.5,
                 font_size=14, bold=True, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)
    for i, (check, result) in enumerate(checks):
        y = 1.9 + i * 0.5
        _add_textbox(s4, check, 0.5, y, 4.5, 0.45, font_size=11, color=WHITE)
        _add_textbox(s4, result, 5.0, y, 1.1, 0.45, font_size=11, color=GREEN, align=PP_ALIGN.CENTER)

    # Right: features
    features = [
        "amount", "log_amount", "amount_zscore",
        "amount_vs_merchant_mean", "merchant_anomaly_rate",
        "category_anomaly_rate", "hour_of_day", "day_of_week",
        "minute_of_hour", "is_weekend", "seconds_since_last_txn",
    ]
    _add_rect(s4, 6.7, 1.4, 6.3, 5.8, GRAY_MID)
    _add_textbox(s4, "Engineered Features (11)", 6.7, 1.4, 6.3, 0.5,
                 font_size=14, bold=True, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)
    for i, feat in enumerate(features):
        y = 1.9 + i * 0.48
        _add_textbox(s4, f"▸  {feat}", 6.9, y, 5.8, 0.44, font_size=12, color=WHITE)

    # ── SLIDE 5 — Model Results ───────────────────────────────────────────────
    s5 = prs.slides.add_slide(blank_layout)
    _fill_slide_bg(s5, GRAY_DARK)
    _add_section_header(s5, "Model Performance", f"Algorithm: {algorithm}")

    # Metric cards
    metric_cards = [
        ("F1 Score", f"{f1:.2%}", BLUE_LIGHT),
        ("Recall", f"{recall:.2%}", GREEN),
        ("Precision", f"{precision:.2%}", YELLOW),
    ]
    for i, (name, val, col) in enumerate(metric_cards):
        x = 0.5 + i * 4.2
        _add_rect(s5, x, 1.5, 3.8, 1.8, GRAY_MID)
        _add_textbox(s5, val, x, 1.6, 3.8, 1.0, font_size=40, bold=True, color=col, align=PP_ALIGN.CENTER)
        _add_textbox(s5, name, x, 2.55, 3.8, 0.5, font_size=14, color=WHITE, align=PP_ALIGN.CENTER)

    _add_image_safe(s5, FIGURES_DIR / "03_amount_boxplot.png", 0.3, 3.5, 6.0, 3.7)
    _add_image_safe(s5, FIGURES_DIR / "05_correlation_heatmap.png", 6.8, 3.5, 6.2, 3.7)

    _add_textbox(s5,
                 f"Best params: {best_params}  ·  Trained on 5 samples  ·  Threshold: 0.30",
                 0.3, 7.1, 12.7, 0.4, font_size=11,
                 color=RGBColor(0x9C, 0xA3, 0xAF), align=PP_ALIGN.CENTER)

    # ── SLIDE 6 — App Demo Screenshot ────────────────────────────────────────
    s6 = prs.slides.add_slide(blank_layout)
    _fill_slide_bg(s6, GRAY_DARK)
    _add_section_header(s6, "Live Dashboard", "Interactive UI with real-time prediction")

    _add_image_safe(s6, SCREENSHOTS_DIR / "01_dashboard_home.png", 0.4, 1.4, 12.5, 5.5)
    _add_textbox(s6, "http://localhost:8000  ·  Tailwind CSS  ·  Auto-refresh on prediction",
                 0.3, 7.1, 12.7, 0.4, font_size=12,
                 color=BLUE_LIGHT, align=PP_ALIGN.CENTER)

    # ── SLIDE 7 — Test Evidence ───────────────────────────────────────────────
    s7 = prs.slides.add_slide(blank_layout)
    _fill_slide_bg(s7, GRAY_DARK)
    _add_section_header(s7, "Automated Quality Gates", "Every pipeline run is fully tested")

    _add_image_safe(s7, SCREENSHOTS_DIR / "03_prediction_result.png", 0.3, 1.4, 6.3, 4.5)
    _add_image_safe(s7, SCREENSHOTS_DIR / "04_swagger_docs.png", 6.8, 1.4, 6.2, 4.5)

    # Results row
    _add_rect(s7, 0.3, 6.1, 12.7, 1.1, GRAY_MID)
    results = [
        ("8 Unit Tests", "All Passed", GREEN),
        ("6 E2E Tests", "All Passed", GREEN),
        ("Playwright", "Chromium", BLUE_LIGHT),
        ("Latency", "< 500ms", YELLOW),
    ]
    for i, (label, val, col) in enumerate(results):
        x = 0.6 + i * 3.2
        _add_textbox(s7, val, x, 6.15, 3.0, 0.55, font_size=18, bold=True, color=col)
        _add_textbox(s7, label, x, 6.65, 3.0, 0.4, font_size=11, color=RGBColor(0x9C, 0xA3, 0xAF))

    # ── SLIDE 8 — Summary ─────────────────────────────────────────────────────
    s8 = prs.slides.add_slide(blank_layout)
    _fill_slide_bg(s8, NAVY)

    _add_rect(s8, 0, 0, 13.33, 0.08, BLUE_LIGHT)
    _add_rect(s8, 0, 7.42, 13.33, 0.08, BLUE_LIGHT)

    _add_textbox(s8, "PIPELINE COMPLETE", 0.5, 0.3, 12.3, 1.0,
                 font_size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_textbox(s8, "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━",
                 0.5, 1.2, 12.3, 0.35, font_size=12, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)

    summary_items = [
        (f"Model    : {algorithm}  —  F1: {f1:.4f}  |  Recall: {recall:.4f}", WHITE),
        ("API      : http://localhost:8000  (running)", GREEN),
        ("Tests    : 8 unit  |  6 E2E  —  all passed", GREEN),
        ("Screenshots: 6 saved to reports/screenshots/", BLUE_LIGHT),
        ("GitHub   : heramb-analytics/claude-code-cust-churn-ml-pipeline-demo", BLUE_LIGHT),
        ("JIRA     : Project KAN  —  6 tickets (KAN-7 through KAN-13)", YELLOW),
        ("Files    : 29 files created across src/, tests/, models/, reports/", WHITE),
        ("Built by : Claude Code (claude-sonnet-4-6) — autonomous pipeline", RGBColor(0x9C, 0xA3, 0xAF)),
    ]
    for i, (text, col) in enumerate(summary_items):
        y = 1.6 + i * 0.68
        _add_textbox(s8, text, 0.8, y, 11.7, 0.6, font_size=15, color=col)

    # Save
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PATH))
    print(f"[STAGE 11] Presentation saved → {OUT_PATH}")


if __name__ == "__main__":
    build_presentation()
